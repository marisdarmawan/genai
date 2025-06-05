import streamlit as st
from openai import OpenAI
import io

# Import library untuk berbagai tipe file
try:
    from PyPDF2 import PdfReader
except ImportError:
    # st.warning("PyPDF2 tidak terinstal. Fitur upload PDF tidak akan berfungsi.")
    PdfReader = None
try:
    from docx import Document as DocxDocument # Alias untuk menghindari konflik nama jika ada
except ImportError:
    # st.warning("python-docx tidak terinstal. Fitur upload DOC/DOCX tidak akan berfungsi.")
    DocxDocument = None
try:
    from pptx import Presentation
except ImportError:
    # st.warning("python-pptx tidak terinstal. Fitur upload PPT/PPTX tidak akan berfungsi.")
    Presentation = None
try:
    import pandas as pd
except ImportError:
    # st.warning("pandas tidak terinstal. Fitur upload XLS/XLSX tidak akan berfungsi.")
    pd = None


# --- App Configuration ---
st.set_page_config(
    page_title="🤖 Chatbot Assistant",
    page_icon="🐑",
    layout="centered",
)

# --- Secrets Management ---
api_key = st.secrets.get("OPENROUTER_API_KEY", None)

# --- OpenAI Client Initialization ---
client = None
if api_key:
    client = OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=api_key,
    )
else:
    st.error("API Key OpenRouter tidak ditemukan. Harap atur di Streamlit Secrets.")


# --- UI Elements ---
st.title("🤖 Chatbot Assistant")
st.caption("Masukkan Pesan yang Anda Inginkan, atau unggah dokumen dari sidebar untuk dibahas.")

if "messages" not in st.session_state:
    st.session_state.messages = [
        {"role": "assistant", "content": "Hai! Aku adalah Chatbot yang siap membantu kamu!"}
    ]
if "uploaded_file_content" not in st.session_state:
    st.session_state.uploaded_file_content = None
if "uploaded_file_name" not in st.session_state:
    st.session_state.uploaded_file_name = None

# Variabel untuk menyimpan hasil file uploader
uploaded_file = None

# --- Sidebar ---
with st.sidebar:
    # --- File Uploader ---
    # Tambahkan tipe file yang didukung ke list
    supported_file_types = ["txt", "md", "pdf"]
    if DocxDocument:
        supported_file_types.extend(["docx", "doc"]) # doc mungkin perlu konversi atau library lain
    if Presentation:
        supported_file_types.append("pptx")
    if pd:
        supported_file_types.extend(["xlsx", "xls"])

    uploaded_file = st.file_uploader( # Pindahkan widget ke sini
        f"Unggah dokumen (opsional, .{', .'.join(supported_file_types)})",
        type=supported_file_types,
        help="Unggah file untuk dibahas dengan chatbot."
    )

    st.markdown("---")
    if st.button("Mulai Percakapan Baru"):
        st.session_state.messages = [
            {"role": "assistant", "content": "Hai! Silahkan ketik pesanmu di sini!"}
        ]
        st.session_state.uploaded_file_content = None
        st.session_state.uploaded_file_name = None
        # Untuk mereset file uploader, kita perlu memanipulasi kunci widgetnya jika ingin
        # membersihkannya secara programmatic setelah tombol ini ditekan.
        # Cara paling mudah adalah dengan st.experimental_rerun() atau st.rerun()
        # yang akan merender ulang seluruh aplikasi, termasuk file uploader ke state awalnya.
        # uploaded_file = None # Ini tidak akan mereset widget file_uploader secara langsung
        st.rerun()

    st.markdown("---")
    st.header("Tentang Bot Ini")
    st.markdown("""
    Bot ini menggunakan model AI dari OpenRouter untuk membantu kebutuhan Anda.
    Cukup ketik permintaanmu dan lihat hasilnya! 🐐🎉
    """)
    st.subheader("Model Digunakan:")
    st.markdown("Deepseek R1 (via OpenRouter)")
    st.markdown("---")
    st.markdown("Dibuat oleh Mohammad Aris Darmawan")


# --- File Processing Logic (tetap di main script, setelah sidebar didefinisikan) ---
if uploaded_file is not None:
    st.session_state.uploaded_file_name = uploaded_file.name
    extracted_text = ""
    file_processed_successfully = False

    try:
        file_content_bytes = uploaded_file.getvalue()
        file_extension = uploaded_file.name.split('.')[-1].lower()

        if file_extension == "txt" or uploaded_file.type == "text/plain" or file_extension == "md" or uploaded_file.type == "text/markdown":
            extracted_text = file_content_bytes.decode("utf-8", errors="replace")
            file_processed_successfully = True
        elif file_extension == "pdf" and PdfReader:
            try:
                pdf_reader = PdfReader(io.BytesIO(file_content_bytes))
                for page in pdf_reader.pages:
                    extracted_text += page.extract_text() + "\n"
                file_processed_successfully = True
            except Exception as e:
                st.error(f"Gagal membaca file PDF: {e}")
        elif (file_extension == "docx" or file_extension == "doc") and DocxDocument:
            if file_extension == "docx":
                try:
                    doc = DocxDocument(io.BytesIO(file_content_bytes))
                    for para in doc.paragraphs:
                        extracted_text += para.text + "\n"
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                extracted_text += cell.text + "\t"
                            extracted_text += "\n"
                    file_processed_successfully = True
                except Exception as e:
                    st.error(f"Gagal membaca file DOCX: {e}")
            else: # .doc
                 st.warning("File .doc mungkin tidak didukung sepenuhnya. Coba konversi ke .docx.")
        elif file_extension == "pptx" and Presentation:
            try:
                prs = Presentation(io.BytesIO(file_content_bytes))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
                file_processed_successfully = True
            except Exception as e:
                st.error(f"Gagal membaca file PPTX: {e}")
        elif (file_extension == "xlsx" or file_extension == "xls") and pd:
            try:
                xls_file = pd.ExcelFile(io.BytesIO(file_content_bytes))
                for sheet_name in xls_file.sheet_names:
                    df = xls_file.parse(sheet_name)
                    extracted_text += f"--- Sheet: {sheet_name} ---\n"
                    extracted_text += df.to_string(index=False) + "\n\n"
                file_processed_successfully = True
            except Exception as e:
                st.error(f"Gagal membaca file Excel: {e}")
        else:
            st.warning(f"Tipe file '{file_extension}' tidak didukung atau library yang dibutuhkan tidak terinstal.")

        if file_processed_successfully and extracted_text.strip(): #
            st.session_state.uploaded_file_content = extracted_text #
            # st.session_state.uploaded_file_name sudah diatur sebelumnya saat file diunggah

            # Pesan konfirmasi tidak lagi ditambahkan ke histori chat.
            # Variabel user_upload_message dan assistant_ack_message bisa dihapus atau dikomentari.
            # Logika untuk menghindari duplikasi (last_user_message, last_assistant_message, dan kondisinya)
            # juga tidak diperlukan lagi untuk pesan-pesan ini.

            # Anda bisa menambahkan notifikasi yang tidak masuk ke histori chat, misalnya:
            st.info(f"File '{st.session_state.uploaded_file_name}' berhasil diproses. Konteksnya akan digunakan untuk pertanyaan Anda selanjutnya.")
            
            # Pastikan tidak ada st.rerun() yang dieksekusi di sini yang mungkin
            # menyebabkan loop atau perilaku tak terduga lainnya terkait file upload.
        
        elif file_processed_successfully and not extracted_text.strip():
            st.info(f"File '{st.session_state.uploaded_file_name}' berhasil diproses tetapi tidak ditemukan konten teks yang bisa diekstrak.")
            st.session_state.uploaded_file_content = None
        elif not file_processed_successfully and uploaded_file:
             st.session_state.uploaded_file_content = None

    except Exception as e:
        st.error(f"Terjadi kesalahan umum saat memproses file '{uploaded_file.name}': {e}")
        st.session_state.uploaded_file_content = None

# Display existing messages
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# --- Chat Input and Logic ---
user_prompt = st.chat_input("Ketik pesanmu untuk chatbot di sini...") #

if user_prompt and client: #
    st.session_state.messages.append({"role": "user", "content": user_prompt}) #
    with st.chat_message("user"): #
        st.markdown(user_prompt) #

    with st.chat_message("assistant"): #
        message_placeholder = st.empty() #
        full_response = "" #
        try:
            # Buat daftar pesan awal dari histori
            payload_messages = [{"role": m["role"], "content": m["content"]} for m in st.session_state.messages] #
            
            # Tambahkan konteks file sebagai pesan sistem jika ada
            if st.session_state.get("uploaded_file_content"):
                context_system_message_content = (
                    f"Anda memiliki akses ke konten dari file bernama '{st.session_state.uploaded_file_name}'. "
                    "Berikut adalah kontennya:\n--- MULAI KONTEKS FILE ---\n"
                    f"{st.session_state.uploaded_file_content}\n"
                    "--- AKHIR KONTEKS FILE ---\n"
                    "Jika pertanyaan pengguna berkaitan dengan informasi dalam file ini, gunakanlah konteks ini untuk menjawab. "
                    "Jangan menyebutkan secara eksplisit bahwa Anda mengambil informasi dari file kecuali jika itu relevan dengan pertanyaan (misalnya, jika pengguna bertanya tentang file itu sendiri)."
                )
                
                # Strategi: Sisipkan pesan sistem ini di awal payload,
                # tapi hindari duplikasi jika pesan sistem yang sama persis sudah ada.
                # Ini akan dikirim setiap kali pengguna mengirim pesan selama file_content ada.
                # Perhatikan batasan token model Anda.
                context_message = {"role": "system", "content": context_system_message_content}
                
                is_context_already_prepended = False
                if payload_messages and payload_messages[0]["role"] == "system" and payload_messages[0]["content"] == context_system_message_content:
                    is_context_already_prepended = True
                
                if not is_context_already_prepended:
                    payload_messages.insert(0, context_message)

            completion = client.chat.completions.create( #
                model="deepseek/deepseek-r1-0528:free", # Ganti dengan model pilihan Anda jika perlu #
                messages=payload_messages,
                stream=True, #
            )
            # ... (sisa kode untuk streaming response)
            for chunk in completion:
                if chunk.choices[0].delta.content is not None:
                    full_response += chunk.choices[0].delta.content
                    message_placeholder.markdown(full_response + "▌") # Karakter penanda ketik
            message_placeholder.markdown(full_response)
        except Exception as e:
            st.error(f"Oops! Terjadi kesalahan: {e}")
            full_response = "Maaf, aku sedang tidak bisa membantumu saat ini. 🥺"
            message_placeholder.markdown(full_response)

    st.session_state.messages.append({"role": "assistant", "content": full_response})

elif not client and user_prompt:
    st.warning("Chatbot tidak aktif karena API Key belum dikonfigurasi dengan benar.")
